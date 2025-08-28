import os
from copy import deepcopy
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData

# ---------- CONFIG ----------
EXCEL_FILE   = "datasheet_imarc.xlsx"
PPT_TEMPLATE = "template.pptx"
PPT_OUT      = "updated_poc.pptx"

TABLE_DECIMALS   = 1
ROW_H_HEADER_IN  = 0.26
ROW_H_BODY_IN    = 0.22
MARGIN_IN        = 0.20

# ---------- HELPERS ----------
def emu_to_in(emu): return float(emu) / float(Emu(914400))
def in_to_emu(inches): return int(inches * 914400)

def musd(x, decimals=TABLE_DECIMALS):
    if pd.isna(x): return ""
    return f"{float(x):,.{decimals}f}"

def unit_label_from_summary(units_str):
    u = (units_str or "").lower()
    if "million" in u: return "Million US$"
    if "billion" in u: return "Billion US$"
    return units_str or ""

def cagr(v0, v1, n_years):
    try:
        v0 = float(v0); v1 = float(v1)
        if v0 <= 0 or v1 <= 0 or n_years <= 0: return ""
        return ((v1 / v0) ** (1.0 / n_years) - 1.0) * 100.0
    except Exception:
        return ""

def fmt_pct(p): return "" if p == "" or pd.isna(p) else f"{float(p):.1f}%"

def safe_set_paragraph(shape, text):
    if not getattr(shape, "has_text_frame", False): return
    tf = shape.text_frame; tf.clear(); tf.paragraphs[0].add_run().text = text

def iter_shapes_recursive(container):
    for sh in container.shapes:
        if sh.shape_type == 6 and hasattr(sh, "shapes"):
            for sub in iter_shapes_recursive(sh): yield sub
        else:
            yield sh

def table_header_text(T):
    header_cells = []
    for c in range(len(T.columns)):
        try: header_cells.append(T.cell(0, c).text.strip())
        except Exception: pass
    return " ".join([h for h in header_cells if h])

def find_table_shape_by_header(slide, keywords):
    kws = [k.lower() for k in keywords]
    for sh in iter_shapes_recursive(slide):
        if getattr(sh, "has_table", False):
            hdr = (table_header_text(sh.table) or "").lower()
            if all(k in hdr for k in kws): return sh, sh.table
    return None, None

def find_text_shape(slide, keywords):
    kws = [k.lower() for k in keywords]
    for sh in iter_shapes_recursive(slide):
        if getattr(sh, "has_text_frame", False):
            if all(k in (sh.text or "").lower() for k in kws): return sh
    return None

# ---------- TABLE STYLING ----------
def style_table_basic(T):
    widths = [Inches(2.4), Inches(1.15), Inches(1.25), Inches(1.25), Inches(1.25)]
    for i in range(min(len(T.columns), len(widths))):
        T.columns[i].width = widths[i]
    for r in range(len(T.rows)):
        T.rows[r].height = Inches(ROW_H_HEADER_IN if r == 0 else ROW_H_BODY_IN)
        for c in range(len(T.columns)):
            tf = T.cell(r, c).text_frame
            tf.word_wrap = False
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.RIGHT
                for run in p.runs: run.font.size = Pt(9)

def rows_with_unit(rows, unit_label):
    out = []
    for r in rows:
        if len(r) == 4: out.append((r[0], unit_label, r[1], r[2], r[3]))
        elif len(r) == 5: out.append(r)
    return out

def ensure_table_body_rows(T, needed_rows):
    have = len(T.rows) - 1
    to_add = max(0, needed_rows - have)
    if to_add == 0: return
    last_tr = T._tbl.tr_lst[-1]
    for _ in range(to_add): T._tbl.append(deepcopy(last_tr))

def fill_table_body(T, rows, unit_label):
    rows = rows_with_unit(rows, unit_label)
    ensure_table_body_rows(T, len(rows))
    for r_idx, row_vals in enumerate(rows, start=1):
        for c_idx, val in enumerate(row_vals[:len(T.columns)]):
            T.cell(r_idx, c_idx).text = str(val)
    style_table_basic(T)

# ---------- CONTINUATION / CLONE ----------

def get_blank_layout(ppt):
    for i, layout in enumerate(ppt.slide_layouts):
        if "blank" in (getattr(layout, "name", "") or "").lower(): return i
    return len(ppt.slide_layouts) - 1

def clone_shape_to_slide(src_shape, dst_slide):
    new_el = deepcopy(src_shape._element)
    dst_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return dst_slide.shapes[-1]

def new_blank_slide(ppt): return ppt.slides.add_slide(ppt.slide_layouts[get_blank_layout(ppt)])

def move_off_slide(shape, ppt):  # visually hide empty bands if needed
    shape.top = ppt.slide_height + in_to_emu(1)

# ---------- ROBUST SHEET PARSER ----------

def read_sales_value_table(xlsx_path, sheet_name):
    raw = pd.read_excel(xlsx_path, sheet_name=sheet_name, header=None)
    title_idx = None
    for i in range(min(len(raw), 80)):
        cell0 = str(raw.iloc[i, 0]).strip() if pd.notna(raw.iloc[i, 0]) else ""
        if ("market breakup" in cell0.lower()) and ("sales value" in cell0.lower()):
            title_idx = i; break
    if title_idx is None:
        raise ValueError(f"Couldn't find 'Sales Value' title in sheet '{sheet_name}'")
    header_idx = title_idx + 1
    header = [str(x).strip() for x in raw.iloc[header_idx].tolist()]
    data_rows = []
    for r in range(header_idx + 1, len(raw)):
        first_text = str(raw.iloc[r, 0]).strip() if pd.notna(raw.iloc[r, 0]) else ""
        if first_text.lower().startswith("market breakup"): break
        if first_text == "" and len(data_rows) > 0: break
        data_rows.append(raw.iloc[r].tolist())
        if first_text.lower() == "total": break
    df = pd.DataFrame(data_rows, columns=header)
    name_col = df.columns[0]
    def pick_year(columns, year):
        for c in columns:
            if str(c).strip() == str(year): return c
        for c in columns:
            try:
                if float(str(c)) == float(year): return c
            except Exception: pass
        return None
    col24 = pick_year(df.columns, 2024); col33 = pick_year(df.columns, 2033)
    if col24 is None or col33 is None:
        raise ValueError(f"{sheet_name}: couldn't find 2024/2033 in {list(df.columns)}")
    bad = {"total", "type", "source", "region", ""}
    mask = df[name_col].astype(str).str.strip().str.lower()
    df = df[~mask.isin(bad)]
    df = df[[name_col, col24, col33]].copy()
    df[col24] = pd.to_numeric(df[col24], errors="coerce")
    df[col33] = pd.to_numeric(df[col33], errors="coerce")
    df = df.dropna(subset=[col24, col33])
    df = df.rename(columns={col24: 2024, col33: 2033})
    return df

def series_from_sheet(df):
    name_col = df.columns[0]
    out = []
    for _, row in df.iterrows():
        name = str(row[name_col]).strip()
        v24 = row[2024]; v33 = row[2033]
        out.append((name, musd(v24), musd(v33), fmt_pct(cagr(v24, v33, 9))))
    return out

# ---------- SPACE CALCS ----------

def max_body_rows_that_fit(top_in, limit_in):
    avail = max(0.0, limit_in - top_in - ROW_H_HEADER_IN)
    return int(avail // ROW_H_BODY_IN)

def expand_on_slide_to_boundary(T_shape, rows, unit_label, boundary_top_emu, pres_height_emu):
    top_in = emu_to_in(T_shape.top)
    boundary_in     = emu_to_in(boundary_top_emu) - MARGIN_IN
    slide_bottom_in = emu_to_in(pres_height_emu) - MARGIN_IN
    limit_in = min(boundary_in, slide_bottom_in)
    max_total = max_body_rows_that_fit(top_in, limit_in)
    show_n = min(max_total, len(rows))
    fill_table_body(T_shape.table, rows[:show_n], unit_label)
    return rows[show_n:]

# ---------- LOAD ----------
summary  = pd.read_excel(EXCEL_FILE, sheet_name="Summary", index_col=0)
forecast = pd.read_excel(EXCEL_FILE, sheet_name="Sales_Forecast")

by_type   = read_sales_value_table(EXCEL_FILE, "By_Type")
by_source = read_sales_value_table(EXCEL_FILE, "By_Source")
by_region = read_sales_value_table(EXCEL_FILE, "By_Region")

market_name = summary.loc["Market Name", "Value"]
units_full  = summary.loc["Units", "Value"]
unit_label  = unit_label_from_summary(units_full)

value_2024     = forecast.loc[forecast["Year"] == 2024, "Sales Value (Million USD)"].values[0]
value_2033     = forecast.loc[forecast["Year"] == 2033, "Sales Value (Million USD)"].values[0]
cagr_2019_2024 = forecast.loc[forecast["Year"] == 2024, "CAGR 2019–2024 (%)"].values[0]
cagr_2024_2033 = cagr(value_2024, value_2033, 9)

rows_type = series_from_sheet(by_type)
rows_src  = series_from_sheet(by_source)
rows_reg  = series_from_sheet(by_region)

ppt = Presentation(PPT_TEMPLATE)
pres_height_emu = ppt.slide_height

# ---------------- SLIDE 34 ----------------
slide34 = ppt.slides[33]
for sh in slide34.shapes:
    if getattr(sh, "has_text_frame", False):
        t = sh.text
        if "Global Food Flavors Market" in t:
            sh.text = t.replace("Global Food Flavors Market", market_name)
        if "The global food flavors market reached" in t:
            sh.text = (
                f"The global {market_name.lower()} market reached a value of "
                f"US$ {value_2024:,.0f} Million in 2024, growing at a CAGR of "
                f"{cagr_2019_2024}% during 2019–2024."
            )
        if "Food Flavors Market" in t:
            sh.text = t.replace("Food Flavors Market", market_name)

forecast_19_24 = forecast[forecast["Year"].between(2019, 2024)]
chart_data = CategoryChartData()
chart_data.categories = forecast_19_24["Year"].tolist()
chart_data.add_series("Sales Value (Million USD)", forecast_19_24["Sales Value (Million USD)"].tolist())
for sh in slide34.shapes:
    if getattr(sh, "has_chart", False):
        sh.chart.replace_data(chart_data)
        break

# Optional AI paragraph (keep/fallback)
ai_paragraph = None
try:
    import google.generativeai as genai
    api_key = ("AIzaSyD4LOjhJtiC2mdRZJYGSa5iavCBRTZfTKU")
    if api_key:
        genai.configure(api_key=api_key)
        prompt = (
            f"Write a professional market analysis paragraph for {market_name}. "
            f"It reached US$ {value_2024:,.0f} Million in 2024 with a CAGR of {cagr_2019_2024}% during 2019–2024. "
            "Highlight key growth drivers, emerging trends, and the 2025–2033 outlook in 2–3 sentences. "
            "Keep it objective and concise."
        )
        ai_paragraph = genai.GenerativeModel("models/gemini-1.5-flash").generate_content(prompt).text.strip()
except Exception:
    ai_paragraph = None

fallback_paragraph = (
    f"The {market_name.lower()} market reached US$ {value_2024:,.0f} Million in 2024 "
    f"after expanding at {cagr_2019_2024}% CAGR during 2019–2024. Demand is supported by shifting consumer "
    "preferences, innovations in production, and broader end-use adoption. "
    "From 2025–2033, continued product innovation and expansion in emerging regions should sustain growth."
)
for sh in slide34.shapes:
    if getattr(sh, "has_text_frame", False) and "Additionally, advancements" in sh.text:
        safe_set_paragraph(sh, ai_paragraph or fallback_paragraph)

# ---------------- SLIDE 33 ----------------
slide33 = ppt.slides[32]

# Particulars
_, T_part = find_table_shape_by_header(slide33, ["particulars"])
if T_part:
    T_part.cell(1, 0).text = market_name
    T_part.cell(1, 1).text = unit_label
    T_part.cell(1, 2).text = musd(value_2024)
    T_part.cell(1, 3).text = musd(value_2033)
    T_part.cell(1, 4).text = fmt_pct(cagr_2024_2033)
    style_table_basic(T_part)

# Locate bands
type_label  = find_text_shape(slide33, ["breakup", "type"])
type_shape,  T_type   = find_table_shape_by_header(slide33, ["breakup", "type"])
form_label  = find_text_shape(slide33, ["breakup", "form"])
form_shape,  T_form   = find_table_shape_by_header(slide33, ["breakup", "form"])
region_label = find_text_shape(slide33, ["breakup", "region"])
region_shape, T_region = find_table_shape_by_header(slide33, ["breakup", "region"])

# 1) TYPE expands first
leftover_type = []
if T_type:
    boundary_for_type = form_shape.top if form_shape else ppt.slide_height
    leftover_type = expand_on_slide_to_boundary(type_shape, rows_type, unit_label,
                                                boundary_for_type, pres_height_emu)

if leftover_type:
    # Continuation: TYPE (cont.) + FORM + REGION all on a new slide
    cont = new_blank_slide(ppt)

    new_type_label = clone_shape_to_slide(type_label, cont) if type_label else None
    if getattr(new_type_label, "has_text_frame", False):
        txt = new_type_label.text
        if "(cont" not in txt.lower(): txt = f"{txt} (cont.)"
        new_type_label.text = txt
        new_type_label.top  = in_to_emu(0.7)

    new_type_tbl = clone_shape_to_slide(type_shape, cont)
    new_type_tbl.top = in_to_emu(1.1)
    T_type_new = new_type_tbl.table
    ensure_table_body_rows(T_type_new, len(leftover_type))
    fill_table_body(T_type_new, leftover_type, unit_label)

    y = emu_to_in(new_type_tbl.top) + ROW_H_HEADER_IN + ROW_H_BODY_IN * (len(T_type_new.rows) - 1) + MARGIN_IN

    if form_shape:
        new_form_label = clone_shape_to_slide(form_label, cont) if form_label else None
        if getattr(new_form_label, "has_text_frame", False):
            txt = new_form_label.text
            if "form" in txt.lower() and "source" not in txt.lower():
                txt = txt.replace("Breakup by Form", "Breakup by Source")
            new_form_label.text = txt
            new_form_label.top  = in_to_emu(y); y += 0.35

        new_form_tbl = clone_shape_to_slide(form_shape, cont)
        new_form_tbl.top = in_to_emu(y)
        fill_table_body(new_form_tbl.table, rows_src, unit_label)
        y = emu_to_in(new_form_tbl.top) + ROW_H_HEADER_IN + ROW_H_BODY_IN * (len(new_form_tbl.table.rows) - 1) + MARGIN_IN

    if region_shape:
        new_reg_label = clone_shape_to_slide(region_label, cont) if region_label else None
        if new_reg_label: new_reg_label.top = in_to_emu(y); y += 0.35
        new_reg_tbl = clone_shape_to_slide(region_shape, cont)
        new_reg_tbl.top = in_to_emu(y)
        fill_table_body(new_reg_tbl.table, rows_reg, unit_label)

    # hide Form/Region bands on original slide (to avoid empty/dup bands)
    if form_shape:  move_off_slide(form_shape, ppt)
    if form_label:  move_off_slide(form_label, ppt)
    if region_shape: move_off_slide(region_shape, ppt)
    if region_label: move_off_slide(region_label, ppt)

else:
    # 2) TYPE fits. Let FORM use remaining vertical space.
    leftover_src = []
    if T_form:
        boundary_for_form = region_shape.top if region_shape else ppt.slide_height
        leftover_src = expand_on_slide_to_boundary(form_shape, rows_src, unit_label,
                                                   boundary_for_form, pres_height_emu)

    if not T_form:
        # no Form band in template: just fill Region and exit
        if T_region: fill_table_body(T_region, rows_reg, unit_label)

    else:
        if len(rows_src) == 0:
            # no rows to show (edge case)
            if T_region: fill_table_body(T_region, rows_reg, unit_label)

        elif leftover_src and T_region:
            # Not all Form rows fit -> move Region to continuation and finish Form
            cont = new_blank_slide(ppt)

            # FORM (cont.) at top
            new_form_label = clone_shape_to_slide(form_label, cont)
            if getattr(new_form_label, "has_text_frame", False):
                txt = new_form_label.text
                if "form" in txt.lower() and "source" not in txt.lower():
                    txt = txt.replace("Breakup by Form", "Breakup by Source")
                if "(cont" not in txt.lower(): txt += " (cont.)"
                new_form_label.text = txt
                new_form_label.top  = in_to_emu(0.7)

            new_form_tbl = clone_shape_to_slide(form_shape, cont)
            new_form_tbl.top = in_to_emu(1.1)
            ensure_table_body_rows(new_form_tbl.table, len(leftover_src))
            fill_table_body(new_form_tbl.table, leftover_src, unit_label)

            # REGION below
            y = emu_to_in(new_form_tbl.top) + ROW_H_HEADER_IN + ROW_H_BODY_IN * (len(new_form_tbl.table.rows) - 1) + MARGIN_IN
            new_reg_label = clone_shape_to_slide(region_label, cont) if region_label else None
            if new_reg_label: new_reg_label.top = in_to_emu(y); y += 0.35
            new_reg_tbl = clone_shape_to_slide(region_shape, cont)
            new_reg_tbl.top = in_to_emu(y)
            fill_table_body(new_reg_tbl.table, rows_reg, unit_label)

            # hide Region band on original slide so you don't see an empty table
            move_off_slide(region_shape, ppt)
            if region_label: move_off_slide(region_label, ppt)

        elif leftover_src and not T_region:
            # If there is no Region band, put Form (cont.) alone on a continuation
            cont = new_blank_slide(ppt)
            new_form_label = clone_shape_to_slide(form_label, cont)
            if getattr(new_form_label, "has_text_frame", False):
                txt = new_form_label.text
                if "form" in txt.lower() and "source" not in txt.lower():
                    txt = txt.replace("Breakup by Form", "Breakup by Source")
                if "(cont" not in txt.lower(): txt += " (cont.)"
                new_form_label.text = txt
                new_form_label.top = in_to_emu(0.7)
            new_form_tbl = clone_shape_to_slide(form_shape, cont)
            new_form_tbl.top = in_to_emu(1.1)
            ensure_table_body_rows(new_form_tbl.table, len(leftover_src))
            fill_table_body(new_form_tbl.table, leftover_src, unit_label)

        else:
            # All Form rows fit; now fill Region in-place
            if T_region: fill_table_body(T_region, rows_reg, unit_label)

# Label tidy-up
for sh in iter_shapes_recursive(slide33):
    if getattr(sh, "has_text_frame", False):
        if "Food Flavors Market" in sh.text:
            sh.text = sh.text.replace("Food Flavors Market", market_name)
        if "Breakup by Form" in sh.text and "By Source" not in sh.text:
            sh.text = sh.text.replace("Breakup by Form", "Breakup by Source")

# ---------- SAVE ----------
ppt.save(PPT_OUT)
print(f"POC PPT generated: {PPT_OUT}")